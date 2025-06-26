from flask import Flask, render_template, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import fitz  # PyMuPDF
import docx
import re
import requests

app = Flask(__name__, static_folder='static')

# === Claude 3 via OpenRouter ===
OPENROUTER_API_KEY = "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" 
CLAUDE_MODEL = "deepseek/deepseek-r1:free"

def call_claude(prompt):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": CLAUDE_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7
    }

    response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=data)

    try:
        result = response.json()
        if 'choices' in result:
            return result['choices'][0]['message']['content']
        elif 'error' in result:
            raise ValueError(f"Claude API Error: {result['error'].get('message', 'Unknown error')}")
        else:
            raise ValueError("Unexpected API response format: " + str(result))
    except Exception as e:
        raise ValueError(f"Failed to parse response: {str(e)} - Raw response: {response.text}")


# === Text Extraction ===
def extract_text(file):
    filename = file.filename
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.txt':
        try:
            text = file.read().decode('utf-8')
        except UnicodeDecodeError:
            text = file.read().decode('latin-1')
    elif ext == '.pdf':
        text = ""
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text()
    elif ext in ['.docx', '.doc']:
        doc = docx.Document(file)
        text = '\n'.join([para.text for para in doc.paragraphs])
    else:
        raise ValueError("Unsupported file type.")
    return text

# === Chunking for prompt safety ===
def split_text(text, max_length=4000):
    paras = text.split("\n")
    chunks, current = [], ""
    for para in paras:
        if len(current) + len(para) < max_length:
            current += para + "\n"
        else:
            chunks.append(current.strip())
            current = para + "\n"
    if current:
        chunks.append(current.strip())
    return chunks

# === Claude slide generator ===
def generate_ppt_content(text):
    prompt = f"""You're a helpful assistant converting content into a PowerPoint.

    1. Analyze this document content.
    2. Generate a short main title (max 10 words).
    3. Create a list of the 6 to 8 most important and non-overlapping section titles like a table of contents. Choose only the most essential topics, even in long documents, to avoid slide overflow.
    4. For each section, generate bullet points as a summary, without forcing a fixed number.
    5. If applicable, extract and preserve key figures like years, monetary values, or statistics.
    6. Add a 'Conclusion' section at the end with a 1–2 sentence summary.

    Return everything in this format:

    ---
    Title: [main_title]

    Contents:
    1. [section_title_1]
    2. [section_title_2]
    ...

    Section: [section_title_1]
    - bullet point 1
    - bullet point 2
    ...

    Section: Conclusion
    - [Conclusion summary here]

    Here is the content:
    {text}
    """

    return call_claude(prompt)


# helper func
def find_layout_by_name(prs, target_name, fallback_index=0):
    for layout in prs.slide_layouts:
        if target_name.lower() == layout.name.lower():
            return layout
    return prs.slide_layouts[fallback_index]


# === PPT Creation ===
def create_slides_from_response(claude_output, original_file_title, selected_template):
    template_path = f"templates/{selected_template}.pptx"
    prs = Presentation(template_path)
    prs.slide_width = Inches(13.33)

    # Layouts based on template
    layout_title = find_layout_by_name(prs, "Main Title")
    layout_contents = find_layout_by_name(prs, "Contents Agenda")
    layout_section_header = find_layout_by_name(prs, "Section Header")
    layout_bullets = find_layout_by_name(prs, "Bullet Content")

    # --- First Slide (Title) ---
    title_match = re.search(r'Title:\s*(.+)', claude_output)
    title_text = title_match.group(1).strip() if title_match else "Document Summary"
    slide_1 = prs.slides.add_slide(layout_title)
    slide_1.shapes.title.text = clean_markdown(title_text)

    # Set subtitle or secondary placeholder if available
    content_placeholder = next((ph for ph in slide_1.placeholders if ph.placeholder_format.idx == 1), None)
    if content_placeholder:
        content_placeholder.text = f"Source: {original_file_title}"

    # --- Contents Slide ---
    contents_slide = prs.slides.add_slide(layout_contents)
    contents_slide.shapes.title.text = "Contents"
    content_placeholder = next((ph for ph in contents_slide.placeholders if ph.placeholder_format.idx == 1), None)
    if content_placeholder:
        content_box = content_placeholder.text_frame
        content_box.clear()
        contents = re.findall(r'^\d+\.\s+(.+)', claude_output, re.MULTILINE)
        for item in contents:
            p = content_box.add_paragraph()
            p.text = clean_markdown(item.strip())
            p.level = 0
            p.font.size = Pt(18)

    # --- Section Slides ---
    sections = re.split(r'Section:\s*(.+)', claude_output)[1:]
    for i in range(0, len(sections), 2):
        heading = sections[i].strip()
        bullets_block = sections[i + 1].strip()
        bullets = [line.strip("-• ") for line in bullets_block.split("\n") if line.strip()]
        if not bullets:
            continue

        # Section Header Slide
        section_slide = prs.slides.add_slide(layout_section_header)
        section_slide.shapes.title.text = clean_markdown(heading)

        # Bullet Content Slide
        content_slide = prs.slides.add_slide(layout_bullets)
        content_slide.shapes.title.text = clean_markdown(heading)

        content_placeholder = next((ph for ph in content_slide.placeholders if ph.placeholder_format.idx == 1), None)
        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            for bullet in bullets:
                para = text_frame.add_paragraph()
                para.text = clean_markdown(bullet)
                para.level = 0
                para.font.size = Pt(18)
                para.space_after = Pt(10)
                para.alignment = PP_ALIGN.LEFT

    prs.save("summary.pptx")
    return "summary.pptx"




def clean_markdown(text):
    return re.sub(r"\*+", "", text)


# === Routes ===
@app.route('/')
def index():
    return render_template('index.html')

@app.route('/', methods=['POST'])
def generate():
    try:
        file = request.files['file']
        if not file:
            return "No file uploaded."

        selected_template = request.form.get('template') or "classic"
        
        text = extract_text(file)
        if not text or len(text.strip()) == 0:
            return "Empty or unreadable file."

        chunks = split_text(text)
        full_text = " ".join(chunks[:5])  # Use first 5 chunks (max ~20K chars)

        claude_output = generate_ppt_content(full_text)
        original_file_title = os.path.splitext(file.filename)[0]
        ppt_path = create_slides_from_response(claude_output, original_file_title, selected_template)
        return send_file(ppt_path, as_attachment=True)

    except Exception as e:
        return f"Error: {str(e)}"
    


if __name__ == "__main__":
    app.run(debug=True)
